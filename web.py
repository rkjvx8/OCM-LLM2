import base64
import uuid
from pptx import Presentation
import pandas as pd
import os
import json
from pathlib import Path
from dotenv import load_dotenv
from flask import Flask, render_template, request, jsonify
from openai import OpenAI

load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    raise RuntimeError("Error: OPENAI_API_KEY is not set. Set it in your environment or in a .env file.")

client = OpenAI(api_key=OPENAI_API_KEY)

UPLOAD_FOLDER = "uploads"
IMAGE_FOLDER = Path("static") / "generated_images"
Path(UPLOAD_FOLDER).mkdir(exist_ok=True)
IMAGE_FOLDER.mkdir(parents=True, exist_ok=True)
ALLOWED_EXTENSIONS = {"txt", "pdf", "md", "json", "ppt", "pptx", "xls", "xlsx"}

def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def load_organization_context() -> str:
    """Load all uploaded organizational context files."""
    context_lines = []
    if Path(UPLOAD_FOLDER).exists():
        for file_path in Path(UPLOAD_FOLDER).glob("*"):
            if file_path.is_file():
                try:
                    content = extract_text_from_file(file_path)
                    if content:
                        context_lines.append(f"[From {file_path.name}]\n{content}")
                except Exception as e:
                    print(f"Error reading {file_path}: {e}")
    return "\n\n".join(context_lines)


def extract_text_from_file(file_path: Path) -> str:
    ext = file_path.suffix.lower()
    if ext in {".txt", ".md", ".json"}:
        try:
            return file_path.read_text(encoding="utf-8").strip()
        except Exception:
            return ""

    if ext in {".ppt", ".pptx"}:
        try:
            presentation = Presentation(str(file_path))
            slides_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slides_text.append(shape.text.strip())
            return "\n".join(slides_text)
        except Exception:
            return ""

    if ext in {".xls", ".xlsx"}:
        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
            table_lines = []
            for sheet_name, df in sheets.items():
                table_lines.append(f"Sheet: {sheet_name}")
                table_lines.append(df.fillna("").astype(str).to_csv(sep="\t", index=False))
            return "\n".join(table_lines)
        except Exception:
            return ""

    return ""

SYSTEM_PROMPT_BASE = (
    "You are a PhD-level Change Management and Training Expert Practitioner. "
    "Answer every request using evidence-based research from industrial/organizational psychology, Prosci, ACMP, and other research-based change management models. "
    "Base recommendations on validated science, peer-reviewed I/O psychology, Prosci methodology, ACMP standards, Kotter, ADKAR, Bridges, Lewin, and similar research-backed frameworks. "
    "Do not invent proprietary methods or use anecdotal opinions. "
    "Write in a scholarly yet practical tone, combining academic rigor with practitioner-ready guidance. "
    "Frame recommendations as action-oriented, outcome-driven advice with clear steps, expected results, and adoption-focused guidance. "
    "Use varied formatting such as bold, underline, and emojis to make key actions and outcomes clear. "
    "Avoid using markdown hashtags after every step. "
    "Include best practices and concrete solutions, and avoid vague statements. "
    "Where relevant, describe the desired change outcomes, implementation actions, and how to measure success. "
    "At the end of every best-practice response, include one or more sources in a dedicated Sources section. "
    "Cite specific research, models, frameworks, or best-practice standards that support the recommendation."
)

app = Flask(__name__)


def create_chat_response(user_prompt: str) -> str:
    org_context = load_organization_context()
    
    system_prompt = SYSTEM_PROMPT_BASE
    if org_context:
        system_prompt += f"\n\n--- ORGANIZATIONAL CONTEXT ---\n{org_context}\n--- END CONTEXT ---\n\nUse this organizational context to tailor your recommendations."
    
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=messages,
        temperature=0.3,
        max_tokens=800,
    )

    choice = response.choices[0]
    message = getattr(choice, "message", None)
    if message is None:
        raise ValueError("No message returned from OpenAI response.")

    if isinstance(message, dict):
        content = message.get("content", "")
    else:
        content = getattr(message, "content", None)
        if content is None and hasattr(message, "__getitem__"):
            content = message["content"]

    return content.strip()


def generate_process_image(user_prompt: str, org_context: str) -> str:
    prompt_text = (
        "Create a simple process diagram for a change management initiative based on the following request and context. "
        "Use clear, actionable steps and label the process flow accordingly. "
        f"Request: {user_prompt}. "
    )
    if org_context:
        prompt_text += f"Organizational context: {org_context}. "

    response = client.images.generate(
        model="gpt-image-1",
        prompt=prompt_text,
        size="1024x1024"
    )

    image_data = response.data[0].b64_json
    image_bytes = base64.b64decode(image_data)
    filename = f"process_{uuid.uuid4().hex}.png"
    file_path = IMAGE_FOLDER / filename
    with open(file_path, "wb") as f:
        f.write(image_bytes)

    return f"/static/generated_images/{filename}"


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/chat", methods=["POST"])
def chat():
    data = request.get_json(force=True)
    prompt = data.get("prompt", "").strip()
    create_image = bool(data.get("create_image", False))
    if not prompt:
        return jsonify({"error": "No prompt provided."}), 400

    try:
        answer = create_chat_response(prompt)
        response_data = {"answer": answer}
        if create_image:
            org_context = load_organization_context()
            image_url = generate_process_image(prompt, org_context)
            response_data["image_url"] = image_url
        return jsonify(response_data)
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


@app.route("/api/upload", methods=["POST"])
def upload():
    if "file" not in request.files:
        return jsonify({"error": "No file provided."}), 400
    
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected."}), 400
    
    if not allowed_file(file.filename):
        return jsonify({"error": f"Allowed file types: {', '.join(ALLOWED_EXTENSIONS)}"}), 400
    
    try:
        file_path = Path(UPLOAD_FOLDER) / file.filename
        file.save(file_path)
        return jsonify({"message": f"File '{file.filename}' uploaded successfully."})
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


@app.route("/api/files", methods=["GET"])
def list_files():
    files = []
    if Path(UPLOAD_FOLDER).exists():
        for file_path in Path(UPLOAD_FOLDER).glob("*"):
            if file_path.is_file():
                files.append(file_path.name)
    return jsonify({"files": files})


@app.route("/api/files/<filename>", methods=["DELETE"])
def delete_file(filename: str):
    file_path = Path(UPLOAD_FOLDER) / filename
    if not file_path.exists():
        return jsonify({"error": "File not found."}), 404
    
    try:
        file_path.unlink()
        return jsonify({"message": f"File '{filename}' deleted successfully."})
    except Exception as exc:
        return jsonify({"error": str(exc)}), 500


if __name__ == "__main__":
    port = int(os.getenv("PORT", 5000))
    debug = os.getenv("FLASK_DEBUG", "false").lower() in {"1", "true", "yes"}
    app.run(host="0.0.0.0", port=port, debug=debug)
